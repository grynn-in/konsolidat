"""Konsolidat — 1-slide overview (bold, minimal design)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette ──
NAVY      = RGBColor(0x0F, 0x0F, 0x1A)
MID_NAVY  = RGBColor(0x1A, 0x1A, 0x2E)
INDIGO    = RGBColor(0x6C, 0x6B, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF6, 0xF1)
GRAY60    = RGBColor(0x88, 0x88, 0x99)
GRAY40    = RGBColor(0x55, 0x55, 0x66)
BODY      = RGBColor(0x33, 0x33, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SL = prs.slides.add_slide(prs.slide_layouts[6])

W = prs.slide_width
H = prs.slide_height


def rect(l, t, w, h, fill=None, line=None):
    s = SL.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.fill.solid(); s.line.color.rgb = line; s.line.width = Pt(1)
    return s


def txt(l, t, w, h, text, sz=12, color=BODY, bold=False, align=PP_ALIGN.LEFT, font='Calibri', spacing=None):
    tb = SL.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:lnSpc'), None)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = font
    return tb


def rich(l, t, w, h, segs, align=PP_ALIGN.LEFT):
    tb = SL.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for text, sz, color, bold, font in segs:
        r = p.add_run()
        r.text = text; r.font.size = Pt(sz)
        r.font.color.rgb = color; r.font.bold = bold
        r.font.name = font or 'Calibri'
    return tb


def multi_para(l, t, w, h, paragraphs, align=PP_ALIGN.LEFT):
    """paragraphs: list of list-of-(text, sz, color, bold, font)"""
    tb = SL.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for pi, segs in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(4) if pi > 0 else Pt(0)
        p.space_after = Pt(0)
        for text, sz, color, bold, font in segs:
            r = p.add_run()
            r.text = text; r.font.size = Pt(sz)
            r.font.color.rgb = color; r.font.bold = bold
            r.font.name = font or 'Calibri'
    return tb


# ═══════════════════════════════════════════════════
# LEFT PANEL — navy, 55% width
# ═══════════════════════════════════════════════════
LEFT_W = Inches(7.2)

rect(Inches(0), Inches(0), LEFT_W, H, NAVY)

# Brand
rich(
    Inches(0.6), Inches(0.4), Inches(6), Inches(0.7),
    [("Konsolid", 40, WHITE, True, 'Calibri'),
     ("at", 40, INDIGO, True, 'Calibri')],
)

# Tagline
txt(Inches(0.6), Inches(1.05), Inches(6), Inches(0.3),
    "Open-Source Enterprise Performance Management",
    sz=12, color=GRAY60)

# Divider line
rect(Inches(0.6), Inches(1.50), Inches(1.2), Pt(2), INDIGO)

# Elevator pitch
rich(
    Inches(0.6), Inches(1.70), Inches(6.0), Inches(0.8),
    [("Multi-entity consolidation, budgeting, cost allocations, and\nvariance analysis \u2014 powered by ", 14, OFF_WHITE, False, 'Calibri'),
     ("=EPM()", 14, INDIGO, True, 'Consolas'),
     (" in the\nspreadsheet your finance team already knows.", 14, OFF_WHITE, False, 'Calibri')],
)

# ── Five capabilities ──
cap_top = Inches(2.85)
caps = [
    ("Excel-Native Reporting",
     "Five worksheet functions query the analytical warehouse directly from Excel. One keystroke refreshes the entire workbook."),
    ("IFRS Consolidation",
     "Foreign exchange translation, cumulative translation adjustments, non-controlling interest splits, and intercompany elimination \u2014 26 data quality assertions per build."),
    ("Cost Allocations",
     "Three-step cascading allocation engine. Distribute costs by headcount, square metres, or revenue with a complete audit trail."),
    ("Budgeting & Variance",
     "Layered budgets with seasonal spread profiles. Actual versus budget variance analysis with favorable/unfavorable logic."),
    ("Open Architecture",
     "MIT-licensed. Built on dbt Core, ClickHouse, and Frappe. Every transformation is a version-controlled SQL model."),
]

y = cap_top
for title, desc in caps:
    # Accent dot
    rect(Inches(0.6), y + Inches(0.06), Inches(0.08), Inches(0.08), INDIGO)

    txt(Inches(0.80), y, Inches(2.8), Inches(0.22),
        title, sz=10, color=INDIGO, bold=True)

    txt(Inches(0.80), y + Inches(0.22), Inches(5.8), Inches(0.45),
        desc, sz=8.5, color=GRAY60)

    y += Inches(0.72)

# ── Footer on left panel ──
txt(Inches(0.6), Inches(6.75), Inches(6), Inches(0.2),
    "konsolid.at   |   Book a 30-minute demo \u2014 no slides, we show you the live stack.",
    sz=9, color=GRAY40)

txt(Inches(0.6), Inches(7.0), Inches(6), Inches(0.2),
    "Airbyte  \u00b7  ClickHouse  \u00b7  dbt Core  \u00b7  Frappe  \u00b7  Excel VBA",
    sz=7, color=GRAY60, font='Consolas')


# ═══════════════════════════════════════════════════
# RIGHT PANEL — light, 45% width
# ═══════════════════════════════════════════════════
RIGHT_X = LEFT_W
RIGHT_W = W - LEFT_W

rect(RIGHT_X, Inches(0), RIGHT_W, H, OFF_WHITE)

rx = RIGHT_X + Inches(0.5)
rw = RIGHT_W - Inches(1.0)

# ── WHY KONSOLIDAT ──
txt(rx, Inches(0.45), rw, Inches(0.3),
    "WHY KONSOLIDAT?",
    sz=13, color=MID_NAVY, bold=True, align=PP_ALIGN.LEFT, font='Calibri')

rect(rx, Inches(0.80), Inches(0.8), Pt(2), INDIGO)

why_data = [
    ("Excel Is the Interface",
     "No new software to learn. Your team builds reports in the tool they\u2019ve mastered over decades. One formula connects Excel directly to the analytical warehouse."),
    ("90% Cost Reduction",
     "No per-user seat fees. No annual licence escalators. Runs on commodity cloud infrastructure. Three-year total cost of ownership under $55,000."),
    ("Transparent & Auditable",
     "44 dbt models and 26 automated tests, all version-controlled. The code is the documentation \u2014 no black-box calculations."),
]

wy = Inches(1.05)
for title, desc in why_data:
    txt(rx, wy, rw, Inches(0.25),
        title, sz=11, color=MID_NAVY, bold=True)
    txt(rx, wy + Inches(0.27), rw, Inches(0.55),
        desc, sz=9, color=BODY)
    wy += Inches(0.90)

# ── COST COMPARISON ──
cost_top = Inches(3.80)
txt(rx, cost_top, rw, Inches(0.3),
    "3-YEAR TOTAL COST OF OWNERSHIP",
    sz=11, color=MID_NAVY, bold=True)

txt(rx, cost_top + Inches(0.25), rw, Inches(0.2),
    "Approximately 50 users",
    sz=8, color=GRAY60)

# Table
tbl_top = cost_top + Inches(0.55)
tbl_shape = SL.shapes.add_table(5, 2, rx, tbl_top, rw, Inches(1.30))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = rw - Inches(2.2)

rows_data = [
    ("", "3-Year Total"),
    ("Tagetik", "$200,000 \u2013 $500,000"),
    ("OneStream", "$300,000 \u2013 $700,000"),
    ("Anaplan", "$700,000 \u2013 $1,400,000"),
    ("Konsolidat", "$20,000 \u2013 $55,000"),
]

for ri, (c0, c1) in enumerate(rows_data):
    is_header = (ri == 0)
    is_hl = (ri == 4)
    for ci, val in enumerate([c0, c1]):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if is_header:
            cell.fill.solid(); cell.fill.fore_color.rgb = MID_NAVY
        elif is_hl:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE8, 0xE7, 0xFF)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE

        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.RIGHT if ci == 1 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(9) if not is_header else Pt(8)
                run.font.bold = is_header or is_hl
                run.font.name = 'Calibri'
                if is_header:
                    run.font.color.rgb = WHITE
                elif is_hl:
                    run.font.color.rgb = INDIGO
                else:
                    run.font.color.rgb = BODY

# Savings banner
banner_top = tbl_top + Inches(1.35)
rect(rx, banner_top, rw, Inches(0.40), INDIGO)
rich(
    rx, banner_top + Inches(0.05), rw, Inches(0.30),
    [("90%+ ", 18, WHITE, True, 'Calibri'),
     ("cost savings versus commercial CPM", 12, WHITE, False, 'Calibri')],
    align=PP_ALIGN.CENTER
)

# ── Stats row ──
stats_top = Inches(6.15)
stats = [("44", "dbt Models"), ("26", "Data Tests"), ("1", "Excel Function"), ("6", "API Endpoints")]
stat_w = rw / 4

for i, (num, label) in enumerate(stats):
    sx = rx + stat_w * i
    txt(sx, stats_top, stat_w, Inches(0.35),
        num, sz=22, color=INDIGO, bold=True, align=PP_ALIGN.CENTER, font='Calibri')
    txt(sx, stats_top + Inches(0.35), stat_w, Inches(0.2),
        label, sz=7, color=GRAY60, bold=True, align=PP_ALIGN.CENTER)

# MIT badge
txt(rx, Inches(6.85), rw, Inches(0.25),
    "MIT Licensed  \u00b7  github.com/grynn-in/konsolidat",
    sz=8, color=GRAY60, align=PP_ALIGN.CENTER, font='Consolas')

prs.save("/home/pd/open_epm/Konsolidat_Overview.pptx")
print("Done")
